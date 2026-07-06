# -*- coding: utf-8 -*-
# One Console AI · ITCEN 신사업 공모전 제안서 v8
# 핵심 메시지: AI Ready Data는 AI Agent 구축의 시작 —
#             One Console AI는 구축→운영→개선 전 생명주기를 관리하는
#             Enterprise AI AgentOps Platform
# 템플릿 문법: 미스터피피티 (표준 헤더·이중 룰·리드 17pt·단일 액센트 #0049F0)
# 실행: py proposal\make_proposal.py → proposal\One_Console_AI_제안서_v8.pptx
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pathlib import Path

# ---------- 팔레트 ----------
BLUE   = RGBColor(0x00, 0x49, 0xF0)
BLUE_D = RGBColor(0x4D, 0x82, 0xFF)
INK    = RGBColor(0x0C, 0x0C, 0x0C)
MUT80  = RGBColor(0x2D, 0x2D, 0x2D)
MUT48  = RGBColor(0x7A, 0x7A, 0x7A)
GRAY   = RGBColor(0xA6, 0xA6, 0xA6)
SILVER = RGBColor(0xBF, 0xBF, 0xBF)
BODY_D = RGBColor(0xCC, 0xCC, 0xCC)
HAIR   = RGBColor(0xD9, 0xD9, 0xD9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PARCH  = RGBColor(0xF5, 0xF5, 0xF7)
PEARL  = RGBColor(0xFB, 0xFB, 0xFB)
TILE   = RGBColor(0x28, 0x29, 0x2B)
TILE2  = RGBColor(0x3A, 0x3A, 0x3E)
BLACK  = RGBColor(0x0C, 0x0C, 0x0C)
FONT   = "Pretendard"

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.62)
CW = SW - ML * 2

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

# ---------- primitives ----------
def _set_font(run, size, bold, color, name=FONT, spc=None):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, name
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', name)
    if spc is not None:
        rPr.set('spc', str(int(spc * 100)))

def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, spacing=1.18):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        if len(para) == 5: p.space_after = Pt(para[4])
        if isinstance(para[0], list):
            for r_ in para[0]:
                run = p.add_run(); run.text = r_[0]
                _set_font(run, r_[1], r_[2], r_[3], spc=(r_[4] if len(r_) > 4 else None))
        else:
            run = p.add_run(); run.text = para[0]
            _set_font(run, para[1], para[2], para[3])
    return tb

def box(slide, x, y, w, h, fill, line=None, r_px=14, lw=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = min(0.5, (r_px / 96) / min(w / 914400, h / 914400))
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp

def pill(slide, x, y, w, h, label, fill, color, size=10.5, bold=True, line=None):
    box(slide, x, y, w, h, fill, line=line, r_px=999)
    text(slide, x, y + Inches(0.01), w, h - Inches(0.02), [(label, size, bold, color)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def conn(slide, x1, y1, x2, y2, color=HAIR, w=1.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    c.shadow.inherit = False
    c.line.color.rgb = color; c.line.width = Pt(w)
    return c

def dot(slide, cx, cy, d, fill, line=None):
    return box(slide, cx - d / 2, cy - d / 2, d, d, fill, line, shape=MSO_SHAPE.OVAL)

def linechart(slide, x, y, w, h, pts, color=BLUE, lw=2.0, dots=True):
    n = len(pts)
    xs = [x + w * i / (n - 1) for i in range(n)]
    ys = [y + h * (1 - p) for p in pts]
    for i in range(n - 1):
        conn(slide, xs[i], ys[i], xs[i + 1], ys[i + 1], color, lw)
    if dots:
        for xi, yi in zip(xs, ys):
            dot(slide, xi, yi, Inches(0.08), color)
    return xs, ys

SHOTS = Path(__file__).resolve().parent / "shots"
def _imgsize(path):
    import struct
    with open(path, "rb") as f:
        head = f.read(26)
    if head[:8] == b"\x89PNG\r\n\x1a\n":
        w, h = struct.unpack(">II", head[16:24])
        return w, h
    return 1600, 1000

def shot_fit(slide, x, y, w, h, name, title="", cap="", bias="top"):
    """브라우저 크롬 + 스크린샷을 w×h 박스에 왜곡 없이 크롭 배치."""
    path = SHOTS / name
    iw, ih = _imgsize(path)
    chrome = Inches(0.32)
    box(slide, x, y, w, chrome, RGBColor(0x33, 0x34, 0x38), r_px=7)
    box(slide, x, y + chrome / 2, w, chrome / 2, RGBColor(0x33, 0x34, 0x38),
        shape=MSO_SHAPE.RECTANGLE)
    for i, c in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFE, 0xBC, 0x2E),
                           RGBColor(0x28, 0xC8, 0x40)]):
        dot(slide, x + Inches(0.18) + Inches(0.17) * i, y + chrome / 2, Inches(0.085), c)
    if title:
        aw = w - Inches(1.5)
        box(slide, x + Inches(0.9), y + Inches(0.07), aw, Inches(0.18),
            RGBColor(0x22, 0x23, 0x26), r_px=99)
        text(slide, x + Inches(0.9), y + Inches(0.07), aw, Inches(0.18),
             [(title, 7.5, False, RGBColor(0x9A, 0x9C, 0xA6))],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    pic = slide.shapes.add_picture(str(path), int(x), int(y + chrome), int(w), int(h))
    box_ar = (w / 914400) / (h / 914400)
    img_ar = iw / ih
    if box_ar >= img_ar:
        crop = (ih - iw / box_ar) / ih
        if bias == "top": pic.crop_bottom = crop
        else: pic.crop_top = pic.crop_bottom = crop / 2
    else:
        crop = (iw - ih * box_ar) / iw
        pic.crop_left = pic.crop_right = crop / 2
    box(slide, x, y + chrome, w, h, None, HAIR, r_px=0, lw=1.0, shape=MSO_SHAPE.RECTANGLE)
    if cap:
        text(slide, x, y + chrome + h + Inches(0.06), w, Inches(0.26),
             [(cap, 9, False, MUT48)], align=PP_ALIGN.CENTER)
    return chrome + h

def thumb(slide, x, y, w, name, label):
    path = SHOTS / name
    iw, ih = _imgsize(path)
    imgh = int(w * ih / iw)
    slide.shapes.add_picture(str(path), int(x), int(y), int(w), int(imgh))
    box(slide, x, y, w, imgh, None, HAIR, r_px=6, lw=1.0)
    text(slide, x, y + imgh + Inches(0.05), w, Inches(0.24),
         [(label, 9.5, True, INK)], align=PP_ALIGN.CENTER)
    return imgh

PAGE = [0]
def base(mode="white"):
    s = prs.slides.add_slide(BLANK)
    bg = {"white": WHITE, "parchment": PARCH, "dark": TILE, "black": BLACK}[mode]
    box(s, 0, 0, SW, SH, bg, shape=MSO_SHAPE.RECTANGLE)
    return s

def content(num, title, lead=None, mode="white"):
    s = base(mode)
    dark = mode in ("dark", "black")
    ink = WHITE if dark else INK
    sub = GRAY if dark else MUT48
    text(s, ML, Inches(0.52), Inches(8.0), Inches(0.26),
         [("Enterprise AI AgentOps Platform - One Console AI", 11, False, sub)])
    text(s, ML, Inches(0.84), Inches(0.6), Inches(0.42), [(num, 24, True, BLUE if not dark else BLUE_D)])
    text(s, ML + Inches(0.66), Inches(0.84), Inches(9.5), Inches(0.42), [(title, 24, True, ink)])
    text(s, SW - Inches(1.6), Inches(0.9), Inches(1.0), Inches(0.3),
         [("ITCEN", 12, True, sub)], align=PP_ALIGN.RIGHT)
    conn(s, ML, Inches(1.37), ML + Inches(0.44), Inches(1.37), BLUE, 2.4)
    conn(s, ML + Inches(0.45), Inches(1.37), SW - ML, Inches(1.37),
         TILE if dark else HAIR, 1.0)
    if lead:
        paras = [(lead, 17, True, ink)] if isinstance(lead, str) else [(rl,) for rl in lead]
        text(s, ML, Inches(1.58), CW, Inches(0.55), paras, spacing=1.22)
    PAGE[0] += 1
    text(s, SW - Inches(1.1), Inches(7.1), Inches(0.5), Inches(0.28),
         [(f"{PAGE[0]:02d}", 9.5, True, GRAY)], align=PP_ALIGN.RIGHT)
    return s

def lead_runs(parts):
    return [[(t, 17, True, BLUE if b else INK) for (t, b) in parts]]

def kw_block(s, x, y, w, kw, body, kw_color=INK):
    text(s, x, y, w, Inches(0.32), [(kw, 15, True, kw_color)])
    conn(s, x, y + Inches(0.42), x + w - Inches(0.25), y + Inches(0.42), HAIR, 1.0)
    text(s, x, y + Inches(0.55), w - Inches(0.15), Inches(0.9), [(body, 11, False, MUT80)],
         spacing=1.3)

def status_pill(s, x, y, done=True, w=Inches(1.6)):
    """STEP 상태 배지: MVP 구현 완료(파랑) / v1.0 로드맵(아웃라인)"""
    if done:
        pill(s, x, y, w, Inches(0.34), "✓ MVP 구현 완료", BLUE, WHITE, 9)
    else:
        pill(s, x, y, w, Inches(0.34), "v1.0 로드맵", WHITE, MUT48, 9, line=HAIR)

def down_arrow(s, cx, y, color=GRAY):
    text(s, cx - Inches(0.15), y, Inches(0.3), Inches(0.3), [("↓", 14, True, color)],
         align=PP_ALIGN.CENTER)

# ============================================================
# 01 · 표지
# ============================================================
s = base("black")
text(s, ML, Inches(0.55), Inches(7), Inches(0.3),
     [("ITCEN 신사업 아이디어 공모전", 12, True, GRAY)])
text(s, SW - Inches(1.72), Inches(0.55), Inches(1.1), Inches(0.3),
     [("ITCEN", 12, True, GRAY)], align=PP_ALIGN.RIGHT)
conn(s, ML, Inches(1.0), ML + Inches(0.44), Inches(1.0), BLUE, 2.4)
conn(s, ML + Inches(0.45), Inches(1.0), SW - ML, Inches(1.0), TILE2, 1.0)
text(s, ML, Inches(1.95), Inches(11.8), Inches(1.9), [
    ([("One Console AI", 54, True, WHITE, -1.2)], 6),
    ([("Enterprise AI ", 22, True, BODY_D), ("AgentOps Platform", 22, True, BLUE_D)],)],
    spacing=1.1)
box(s, ML, Inches(4.0), Inches(0.56), Inches(0.05), BLUE, r_px=99)
text(s, ML, Inches(4.3), Inches(11.5), Inches(1.3), [
    ([("AI Ready Data", 17, True, BLUE_D), ("는 AI Agent 구축의 시작입니다.", 17, True, WHITE)], 5),
    ("One Console AI는 구축부터 운영·개선까지, AI Agent의 전 생명주기를 하나의 플랫폼에서 관리합니다.",
     13, False, GRAY)])
# 생명주기 미니 스트립
lc = ["AI Ready Data", "Knowledge", "Agent 생성", "AgentOps", "Continuous Ops"]
for i, t_ in enumerate(lc):
    x = ML + Inches(2.32) * i
    pill(s, x, Inches(5.5), Inches(2.1), Inches(0.42), t_, TILE, BODY_D if i else BLUE_D, 9.5)
    if i < 4:
        text(s, x + Inches(2.1), Inches(5.53), Inches(0.24), Inches(0.34),
             [("→", 11, True, MUT48)], align=PP_ALIGN.CENTER)
box(s, 0, Inches(6.35), SW, Inches(0.62), TILE, shape=MSO_SHAPE.RECTANGLE)
text(s, ML, Inches(6.5), Inches(12), Inches(0.32),
     [("2026. 07   |   아이디어 제안서   |   제안팀 ____________", 11, False, BODY_D)])

# ============================================================
# 02 · 제안 요약 (00) — One Message 먼저
# ============================================================
s = content("00", "제안 요약",
            lead=lead_runs([("One Message", True), (" — 그리고 네 가지 심사 기준에 대한 답", False)]))
box(s, ML, Inches(2.25), CW, Inches(0.95), TILE)
text(s, ML + Inches(0.45), Inches(2.44), CW - Inches(0.9), Inches(0.62), [
    ([("One Console AI는 AI Ready Data를 시작으로 AI Agent 구축·운영·개선까지 ", 13.5, True, WHITE),
      ("전 생명주기를 하나의 플랫폼에서 관리", 13.5, True, BLUE_D),
      ("하는 Enterprise AI AgentOps Platform입니다.", 13.5, True, WHITE)],)], spacing=1.3)
cards = [
    ("35", "문제 발굴 · Industry Insight", "구축보다 어려운 것은 지속적인 운영",
     "AI Agent를 만드는 것이 문제가 아니라, 지속 운영할 플랫폼이 없다는 것 — 현장에서 반복 확인"),
    ("30", "솔루션 유효성 · 차별성", "AI Ready Data부터 AgentOps까지 5 STEP 통합",
     "Assessment → Data Studio → Knowledge Intelligence → AgentOps → Continuous Operation"),
    ("20", "사업화 가능성", "SI → SaaS → MSP, 반복 매출 구조",
     "동작하는 MVP 기반 즉시 사업화 — 6개월 그룹사 Pilot → 36개월 고객 20~30사"),
    ("15", "실행력", "슬라이드가 아니라 실제 제품으로 제안",
     "콘솔 8개 화면 + REST API 19개 구현 완료 — 발표장 라이브 데모"),
]
gw, gh = Inches(5.9), Inches(1.62)
for i, (score, crit, head, body) in enumerate(cards):
    x = ML + (gw + Inches(0.29)) * (i % 2)
    y = Inches(3.45) + (gh + Inches(0.22)) * (i // 2)
    box(s, x, y, gw, gh, PEARL, HAIR)
    text(s, x + Inches(0.3), y + Inches(0.22), Inches(1.05), Inches(1.2), [
        ([(score, 30, True, BLUE, -0.7)], 0), ("점", 9.5, False, MUT48)])
    text(s, x + Inches(1.35), y + Inches(0.22), gw - Inches(1.7), gh - Inches(0.4), [
        (crit, 9, True, MUT48, 3), (head, 12.5, True, INK, 4), (body, 9.5, False, MUT80)],
        spacing=1.18)

# ============================================================
# 03 · 문제 1/2 (01) — 기준 부재 + 반복 작업
# ============================================================
s = content("01", "현장 문제 발굴 & 정의",
            lead=lead_runs([("AI Agent 구축보다 더 어려운 것은, ", False), ("지속적인 운영", True), ("입니다.", False)]))
text(s, ML, Inches(2.25), CW, Inches(0.35),
     [("도입 의지는 높지만 — 세 가지 질문에 답할 기준과 체계가 없었습니다", 12.5, True, INK)])
qs = [("어디부터?", "어떤 업무부터 AI를\n적용해야 하는가"),
      ("데이터는?", "AI가 활용할 데이터를\n어떻게 준비해야 하는가"),
      ("그 다음은?", "구축 이후 어떻게 운영하고\n지속적으로 개선하는가")]
for i, (q, b_) in enumerate(qs):
    x = ML + Inches(4.1) * i
    box(s, x, Inches(2.7), Inches(3.9), Inches(1.3), PEARL, HAIR)
    text(s, x + Inches(0.3), Inches(2.9), Inches(3.3), Inches(0.95), [
        ([(q, 15, True, BLUE)], 4), (b_, 10.5, False, MUT80)], spacing=1.2)
text(s, ML, Inches(4.25), CW, Inches(0.35),
     [("그 결과 — AI 프로젝트를 수행할 때마다 같은 작업이 반복되었습니다", 12.5, True, INK)])
reps = [("문서 전처리 · AI Ready Data 재구축", "이전 프로젝트의 파이프라인은 소멸"),
        ("프롬프트 · 검색 정책 재설계", "위키·메모·머릿속에 산재"),
        ("Agent 설정 · Workflow 재구성", "봇마다 처음부터 다시"),
        ("운영 로그 · 개선 이력 단절", "무엇이 왜 좋아졌는지 추적 불가")]
for i, (h_, b_) in enumerate(reps):
    x = ML + Inches(6.15) * (i % 2)
    y = Inches(4.7) + Inches(0.95) * (i // 2)
    box(s, x, y, Inches(5.9), Inches(0.82), WHITE, HAIR)
    text(s, x + Inches(0.3), y + Inches(0.13), Inches(5.35), Inches(0.58), [
        ([("✕  ", 12, True, GRAY), (h_, 12, True, INK), ("   " + b_, 9.5, False, MUT48)],)],
        anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 04 · 문제 2/2 (01) — PoC 종료 구조 + Business Insight
# ============================================================
s = content("01", "현장 문제 발굴 & 정의",
            lead=lead_runs([("PoC는 성공해도, ", False), ("확산은 실패한다", True),
                            (" — 운영 체계 부재의 구조적 결과", False)]))
box(s, ML, Inches(2.25), Inches(6.55), Inches(3.55), PEARL, HAIR)
text(s, ML + Inches(0.38), Inches(2.45), Inches(5.8), Inches(0.3),
     [("프로젝트를 거듭해도 — 비용은 다시 들고, 자산은 남지 않는다", 12, True, INK)])
gx, gy, gw_, gh_ = ML + Inches(0.5), Inches(2.9), Inches(5.7), Inches(2.4)
conn(s, gx, gy + gh_, gx + gw_, gy + gh_, HAIR, 1.0)
conn(s, gx, gy, gx, gy + gh_, HAIR, 1.0)
linechart(s, gx + Inches(0.2), gy + Inches(0.2), gw_ - Inches(0.5), gh_ - Inches(0.45),
          [0.18, 0.38, 0.58, 0.78, 0.98], color=SILVER)
linechart(s, gx + Inches(0.2), gy + Inches(0.2), gw_ - Inches(0.5), gh_ - Inches(0.45),
          [0.06, 0.07, 0.06, 0.07, 0.06], color=BLUE)
text(s, gx + gw_ - Inches(1.8), gy + Inches(0.05), Inches(1.9), Inches(0.3),
     [("누적 투입 비용", 9.5, True, MUT48)])
text(s, gx + gw_ - Inches(1.9), gy + gh_ - Inches(0.6), Inches(2.0), Inches(0.3),
     [("남는 조직 자산 ≈ 0", 9.5, True, BLUE)])
for i, l in enumerate(["PoC 1", "2", "3", "4", "… N"]):
    text(s, gx + Inches(0.05) + (gw_ - Inches(0.5)) * i / 4 - Inches(0.3), gy + gh_ + Inches(0.07),
         Inches(0.8), Inches(0.24), [(l, 8.5, False, GRAY)], align=PP_ALIGN.CENTER)
stats = [("8주+", "Agent 구축 리드타임", "AI Ready Data부터 매번 재구축"),
         ("30~40%", "유사 기능 중복 개발", "부서별로 같은 봇을 각자 구축"),
         ("0", "운영 → 개선 사이클", "PoC 종료와 함께 이력도 종료")]
for i, (v, l, sub) in enumerate(stats):
    y = Inches(2.35) + Inches(1.18) * i
    text(s, ML + Inches(6.95), y, Inches(4.9), Inches(1.1), [
        ([(v, 27, True, INK, -0.6), ("   " + l, 12, True, INK)], 3), (sub, 9.5, False, MUT48)])
    if i < 2: conn(s, ML + Inches(6.95), y + Inches(1.02), SW - ML, y + Inches(1.02), HAIR, 1.0)
box(s, ML, Inches(6.0), CW, Inches(0.92), TILE)
text(s, ML + Inches(0.45), Inches(6.17), CW - Inches(0.9), Inches(0.62), [
    ([("Business Insight — ", 13.5, True, BLUE_D),
      ("문제는 AI Agent를 만드는 것이 아니라, AI Agent를 지속적으로 운영할 수 있는 플랫폼이 없다는 것입니다.",
       13.5, True, WHITE)],)], spacing=1.25)
text(s, ML, Inches(7.05), Inches(8), Inches(0.26),
     [("※ 수치는 당사 프로젝트 수행 경험 기반 추정치", 8.5, False, GRAY)])

# ============================================================
# 05 · 솔루션 개요 (02) — 5 STEP
# ============================================================
s = content("02", "Industry Insight 기반 솔루션",
            lead=lead_runs([("AI Ready Data부터 AgentOps까지, ", False), ("하나의 플랫폼으로 통합", True), (".", False)]))
steps5 = [
    ("STEP 1", "AI Ready\nAssessment", ["적용 가능 업무 진단", "데이터 준비 수준 분석", "도입 우선순위 제안"], False),
    ("STEP 2", "AI Ready\nData Studio", ["문서 전처리 자동화", "Metadata · Quality 관리", "Vector DB 구축"], True),
    ("STEP 3", "Knowledge\nIntelligence", ["Ontology 자동 생성", "GraphRAG · Hybrid Search", "관계 기반 지식 관리"], False),
    ("STEP 4", "AgentOps", ["Prompt · Workflow 관리", "Search Policy 관리", "Version · Rollback · Diff"], True),
    ("STEP 5", "Continuous\nAI Operation", ["Feedback · Evaluation", "Monitoring · 로그 분석", "AI 자동 개선"], True),
]
cw5 = Inches(2.29)
for i, (st, h_, bullets, done) in enumerate(steps5):
    x = ML + (cw5 + Inches(0.16)) * i
    box(s, x, Inches(2.3), cw5, Inches(3.15), PEARL if not done else WHITE,
        HAIR if not done else BLUE, lw=1.0 if not done else 1.4)
    text(s, x + Inches(0.2), Inches(2.5), cw5 - Inches(0.4), Inches(0.25),
         [(st, 9.5, True, BLUE)])
    text(s, x + Inches(0.2), Inches(2.76), cw5 - Inches(0.4), Inches(0.62),
         [(h_, 13, True, INK)], spacing=1.05)
    yy = Inches(3.55)
    for b_ in bullets:
        text(s, x + Inches(0.2), yy, cw5 - Inches(0.34), Inches(0.5),
             [("· " + b_, 8.8, False, MUT80)], spacing=1.12)
        yy += Inches(0.48)
    status_pill(s, x + Inches(0.2), Inches(5.0), done, w=cw5 - Inches(0.4))
    if i < 4:
        text(s, x + cw5 - Inches(0.02), Inches(3.6), Inches(0.24), Inches(0.4),
             [("→", 12, True, GRAY)], align=PP_ALIGN.CENTER)
box(s, ML, Inches(5.85), CW, Inches(0.92), TILE)
text(s, ML + Inches(0.45), Inches(6.02), CW - Inches(0.9), Inches(0.62), [
    ([("핵심 메시지 — ", 13.5, True, BLUE_D),
      ("AI Ready Data를 '만드는' 플랫폼이 아니라, AI Agent가 지속적으로 '성장하는' 운영 플랫폼입니다.",
       13.5, True, WHITE)],)], spacing=1.25)
text(s, ML, Inches(6.95), CW, Inches(0.26),
     [("파란 테두리 = 현재 MVP로 동작 검증 완료 · 회색 = v1.0 제품화 로드맵", 8.5, False, GRAY)])

# ============================================================
# 06 · STEP 1·2 (02) — Assessment + Data Studio
# ============================================================
s = content("02", "STEP 1·2 — AI Ready Data가 만들어지는 곳",
            lead=lead_runs([("좋은 Agent는 ", False), ("AI Ready Data", True), ("에서 시작됩니다.", False)]))
box(s, ML, Inches(2.3), Inches(5.7), Inches(1.9), PEARL, HAIR)
text(s, ML + Inches(0.32), Inches(2.5), Inches(5.1), Inches(0.3),
     [([("STEP 1  ", 11, True, BLUE), ("AI Ready Assessment", 13.5, True, INK)],)])
text(s, ML + Inches(0.32), Inches(2.95), Inches(5.1), Inches(1.1), [
    ("· 어떤 업무부터 적용할지 — AI 적용 가능 업무 진단", 10.5, False, MUT80, 5),
    ("· 데이터 준비 수준 분석 → AI 도입 우선순위 제안", 10.5, False, MUT80, 5),
    ("· 도입 컨설팅이 곧 구축·구독 수주의 입구", 10.5, True, BLUE)], spacing=1.2)
status_pill(s, ML + Inches(0.32), Inches(3.85) + Inches(0.0), False)
box(s, ML, Inches(4.4), Inches(5.7), Inches(2.25), WHITE, BLUE, lw=1.4)
text(s, ML + Inches(0.32), Inches(4.6), Inches(5.1), Inches(0.3),
     [([("STEP 2  ", 11, True, BLUE), ("AI Ready Data Studio", 13.5, True, INK)],)])
text(s, ML + Inches(0.32), Inches(5.05), Inches(5.15), Inches(1.3), [
    ("· 이미지·표·비정형 문서 → 추출·전처리·가공 자동화", 10.5, False, MUT80, 5),
    ("· 한국어 출력 규칙 강제 (서론·추측 금지, 사실 위주)", 10.5, False, MUT80, 5),
    ("· Metadata 자동 생성 · Data Quality 관리 (v1.0)", 10.5, False, MUT80, 5),
    ("· 청킹·임베딩 → Vector DB 구축·동기화", 10.5, False, MUT80)], spacing=1.2)
status_pill(s, ML + Inches(0.32), Inches(6.22), True)
shot_fit(s, ML + Inches(6.0), Inches(2.5), Inches(5.5), Inches(3.6), "knowledge.png",
         title="One Console AI · RAG 데이터 파이프라인",
         cap="MVP 구현 화면 — 수집→추출→전처리→가공→적재 파이프라인이 실제로 동작합니다")

# ============================================================
# 07 · STEP 3 (02) — Knowledge Intelligence
# ============================================================
s = content("02", "STEP 3 — Knowledge Intelligence",
            lead=lead_runs([("Vector를 넘어 ", False), ("Ontology · GraphRAG", True),
                            ("로 — 관계를 아는 지식.", False)]))
# 좌: 트리플 추출 플로우
box(s, ML, Inches(2.3), Inches(6.3), Inches(2.6), PEARL, HAIR)
text(s, ML + Inches(0.32), Inches(2.5), Inches(5.7), Inches(0.3),
     [("가공 단계에서 Ontology가 '자동으로' 만들어집니다", 12.5, True, INK)])
flow3 = ["문서 가공", "Triple 추출\n(개체-관계-개체)", "지식그래프 적재", "Ontology 스키마\n확장 제안 → 승인"]
fx = ML + Inches(0.3)
for i, t_ in enumerate(flow3):
    fw3 = Inches(1.38)
    box(s, fx, Inches(3.0), fw3, Inches(0.95), WHITE, HAIR if i != 1 else BLUE, lw=1.2)
    text(s, fx, Inches(3.06), fw3, Inches(0.84), [(t_, 8.8, i == 1, INK if i == 1 else MUT80)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)
    if i < 3:
        text(s, fx + fw3, Inches(3.32), Inches(0.16), Inches(0.3), [("→", 10, True, GRAY)],
             align=PP_ALIGN.CENTER)
    fx += fw3 + Inches(0.16)
text(s, ML + Inches(0.32), Inches(4.15), Inches(5.7), Inches(0.6), [
    ([("예시  ", 9.5, True, BLUE),
      ("X-200 —호환— 부품B — 적용규정— 안전규정 §12  →  \"X-200 호환 부품의 안전 규정은?\" 같은 다단계 질문에 응답",
       9.5, False, MUT48)],)], spacing=1.25)
# 우: Hybrid Search + 효과
box(s, ML + Inches(6.6), Inches(2.3), Inches(5.5), Inches(2.6), WHITE, HAIR)
text(s, ML + Inches(6.92), Inches(2.5), Inches(4.9), Inches(2.2), [
    ("Hybrid Search = Vector + Graph", 12.5, True, INK, 7),
    ("· 의미 유사(Vector) + 관계 탐색(Graph) 병행", 10.5, False, MUT80, 5),
    ("· 답변 근거 경로 설명 가능 — 공공·금융 감사 대응", 10.5, False, MUT80, 5),
    ("· 규정 개정 시 영향받는 문서·봇 자동 식별", 10.5, False, MUT80, 5),
    ("· 상충 지식(모순) 감지 → 개선 큐 이관", 10.5, False, MUT80)], spacing=1.22)
kw_block(s, ML, Inches(5.25), Inches(5.9), "왜 차별점인가",
         "해외 GraphRAG 도구도 'Ontology의 형상관리(버전·Diff·승인)'까지는 하지 않습니다 — 프롬프트·검색·온톨로지를 한 콘솔에서 관리하는 것이 One Console AI의 자리입니다.")
status_pill(s, ML + Inches(6.6), Inches(5.35), False, w=Inches(2.2))
text(s, ML + Inches(6.6), Inches(5.85), Inches(5.4), Inches(0.6),
     [("v1.0 하이브리드 검색 → 확장 단계 Ontology 자동 구축·영향 분석 (로드맵 연동)", 9.5, False, MUT48)],
     spacing=1.25)

# ============================================================
# 08 · STEP 4 (02) — AgentOps
# ============================================================
s = content("02", "STEP 4 — AgentOps",
            lead=lead_runs([("Prompt · Workflow · 검색 정책을 ", False), ("코드처럼 형상관리", True), (".", False)]))
box(s, ML, Inches(2.3), Inches(5.7), Inches(2.0), TILE)
text(s, ML + Inches(0.32), Inches(2.55), Inches(5.1), Inches(1.5), [
    ("하나의 KEY — bot_type", 11, True, BLUE_D, 6),
    ("에이전트 생성 시 정한 키에 Prompt·검색 정책·Workflow가 매달리고, 외부 호출도 이 키 하나로 식별됩니다.",
     10.5, False, BODY_D)], spacing=1.28)
feats = [("Version 관리", "저장할 때마다 새 버전 — 변경 이력 전 추적"),
         ("Rollback", "문제 발생 시 과거 버전으로 원클릭 복귀"),
         ("Diff 관리", "무엇이 왜 바뀌었는지 버전 간 비교"),
         ("Search Policy", "Index·Top-K·유사도·출처 인용 강제 중앙 관리")]
for i, (h_, b_) in enumerate(feats):
    y = Inches(4.5) + Inches(0.56) * i
    text(s, ML, y, Inches(5.7), Inches(0.5), [
        ([("✓  ", 11, True, BLUE), (h_, 11.5, True, INK), ("   " + b_, 9.5, False, MUT48)],)])
status_pill(s, ML, Inches(6.75) - Inches(0.06), True)
shot_fit(s, ML + Inches(6.0), Inches(2.5), Inches(5.5), Inches(3.6), "prompts.png",
         title="One Console AI · System Prompt 형상관리",
         cap="MVP 구현 화면 — 버전 히스토리·롤백·즉시 배포가 실제로 동작합니다")

# ============================================================
# 09 · STEP 5 (02) — Continuous AI Operation
# ============================================================
s = content("02", "STEP 5 — Continuous AI Operation",
            lead=lead_runs([("운영 신호가 되돌아와, ", False), ("AI가 스스로 좋아집니다", True), (".", False)]))
loop5 = [("Feedback", "답변 평가·코멘트 회귀"), ("Evaluation", "★2 이하 → 개선 큐"),
         ("Monitoring", "로그·사용량·멈춘 작업(4h+)"), ("자동 개선", "Prompt 새 버전 → 배포")]
nw = Inches(2.8)
for i, (h_, b_) in enumerate(loop5):
    x = ML + (nw + Inches(0.3)) * i
    box(s, x, Inches(2.3), nw, Inches(1.15), PEARL if i < 3 else WHITE,
        HAIR if i < 3 else BLUE, lw=1.0 if i < 3 else 1.4)
    text(s, x + Inches(0.24), Inches(2.48), nw - Inches(0.46), Inches(0.8), [
        (h_, 12.5, True, BLUE if i == 3 else INK, 4), (b_, 9.5, False, MUT80)], spacing=1.15)
    if i < 3:
        text(s, x + nw + Inches(0.02), Inches(2.62), Inches(0.28), Inches(0.4),
             [("→", 13, True, GRAY)], align=PP_ALIGN.CENTER)
box(s, ML, Inches(3.8), Inches(5.7), Inches(2.3), PEARL, HAIR)
text(s, ML + Inches(0.35), Inches(3.98), Inches(5.0), Inches(0.3),
     [("개선 사이클 반복 → 만족도 상승 (운영 예시)", 11.5, True, INK)])
gx, gy, gw_, gh_ = ML + Inches(0.5), Inches(4.42), Inches(4.75), Inches(1.35)
conn(s, gx, gy + gh_, gx + gw_, gy + gh_, HAIR, 1.0)
xs, ys = linechart(s, gx + Inches(0.15), gy + Inches(0.1), gw_ - Inches(0.4), gh_ - Inches(0.3),
                   [0.15, 0.3, 0.42, 0.58, 0.75, 0.9], color=BLUE)
text(s, xs[0] - Inches(0.18), ys[0] + Inches(0.08), Inches(0.7), Inches(0.24), [("86%", 9, True, GRAY)])
text(s, xs[-1] - Inches(0.5), ys[-1] - Inches(0.28), Inches(0.8), Inches(0.24), [("94%", 10, True, BLUE)])
status_pill(s, ML + Inches(0.35), Inches(5.55) + Inches(0.16), True)
shot_fit(s, ML + Inches(6.0), Inches(3.6), Inches(5.5), Inches(2.8), "operations.png",
         title="One Console AI · 운영 대시보드",
         cap="MVP 구현 화면 — 로그·피드백·멈춘 작업 정리가 실제로 동작합니다")
text(s, ML, Inches(6.55), Inches(5.7), Inches(0.5),
     [("개선 전후가 버전으로 비교·증명 — \"운영할수록 좋아진다\"가 데이터로 남습니다.", 11, True, INK)],
     spacing=1.25)

# ============================================================
# 10 · 제품 둘러보기 (02)
# ============================================================
s = content("02", "제품 둘러보기",
            lead=lead_runs([("슬라이드가 아니라, ", False), ("실제 동작하는 제품", True), ("입니다.", False)]))
shot_fit(s, ML, Inches(2.4), Inches(6.7), Inches(3.9), "dashboard.png",
         title="One Console AI · localhost:8000",
         cap="실데이터로 동작하는 대시보드 — 생명주기 전 단계를 한 화면에")
rx = ML + Inches(7.05)
text(s, rx, Inches(2.4), Inches(4.5), Inches(0.3), [("함께 구현된 핵심 화면", 12, True, INK)])
grid = [("agents.png", "Agent Config"), ("search.png", "Search Policy"),
        ("improve.png", "지속 개선 큐"), ("api.png", "ITCEN-API · Chat")]
tw = Inches(2.14)
for i, (nm, lb) in enumerate(grid):
    tx = rx + (tw + Inches(0.24)) * (i % 2)
    ty = Inches(2.9) + Inches(1.78) * (i // 2)
    thumb(s, tx, ty, tw, nm, lb)
pill(s, rx, Inches(6.28), Inches(4.5), Inches(0.44),
     "8개 화면 · 19개 REST API 모두 실제 구현 완료", BLUE, WHITE, 10.5)

# ============================================================
# 11 · 개발 연동 (02) — ITCEN-API
# ============================================================
s = content("02", "개발 연동 · ITCEN-API",
            lead=lead_runs([("등록한 형상을, ", False), ("REST로 그대로", True), (" 쓴다.", False)]),
            mode="parchment")
box(s, ML, Inches(2.3), Inches(6.35), Inches(3.15), INK)
code = [
    ("# 형상 가져오기", GRAY),
    ("GET  /agent/prompt/{bot_type}", BLUE_D),
    ("GET  /agent/rag-config/{bot_type}", BLUE_D),
    ("", WHITE),
    ("# 대화 호출 — 형상(Prompt+검색) 자동 적용", GRAY),
    ("POST /agent/chat", WHITE),
    ('     {"bot_type": "hr-assistant",', BODY_D),
    ('      "message": "육아휴직은 최대 얼마나?"}', BODY_D),
    ("", WHITE),
    ("# 피드백 회귀 — Continuous Operation의 시작", GRAY),
    ("POST /agent/feedback", WHITE),
]
tb = text(s, ML + Inches(0.4), Inches(2.56), Inches(5.6), Inches(2.7),
          [(t_, 10.5, False, c_) for (t_, c_) in code], spacing=1.26)
for p in tb.text_frame.paragraphs:
    for r in p.runs: r.font.name = "Consolas"
chips = [("19", "엔드포인트"), ("3", "형상 API"), ("0초", "저장=배포")]
for i, (v, l) in enumerate(chips):
    x = ML + Inches(2.17) * i
    box(s, x, Inches(5.65), Inches(1.98), Inches(0.72), WHITE, HAIR)
    text(s, x, Inches(5.75), Inches(1.98), Inches(0.55), [
        ([(v, 15, True, BLUE, -0.3), ("  " + l, 9.5, False, MUT48)],)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rt = ML + Inches(6.7)
text(s, rt, Inches(2.3), Inches(4.9), Inches(0.3),
     [("자동 생성 API 문서 · Swagger UI", 12, True, INK)])
shot_fit(s, rt, Inches(2.72), Inches(4.9), Inches(3.3), "swagger.png",
         title="localhost:8000/docs")
pill(s, rt, Inches(6.5), Inches(4.9), Inches(0.44),
     "LIVE DEMO — 발표장에서 실제 제품으로 시연", BLUE, WHITE, 10.5)

# ============================================================
# 12 · 고객 가치 (03) — 기존 vs One Console 루프
# ============================================================
s = content("03", "고객 가치 창출",
            lead=lead_runs([("운영할수록 ", False), ("더 똑똑해지는 AI", True),
                            (" — 변경이 자동으로 반영되는 구조", False)]))
# 좌: 기존
box(s, ML, Inches(2.3), Inches(3.35), Inches(4.0), PEARL, HAIR)
text(s, ML + Inches(0.28), Inches(2.48), Inches(2.8), Inches(0.3), [("기존 방식", 12, True, MUT48)])
old = ["데이터 변경", "재개발", "재배포", "…반복"]
for i, t_ in enumerate(old):
    y = Inches(2.95) + Inches(0.82) * i
    box(s, ML + Inches(0.4), y, Inches(2.55), Inches(0.56), WHITE, HAIR)
    text(s, ML + Inches(0.4), y + Inches(0.02), Inches(2.55), Inches(0.5),
         [(t_, 11, True, MUT48)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3: down_arrow(s, ML + Inches(1.68), y + Inches(0.56))
# 우: One Console 자동 루프 (6단계 압축)
box(s, ML + Inches(3.65), Inches(2.3), Inches(4.35), Inches(4.0), WHITE, BLUE, lw=1.4)
text(s, ML + Inches(3.95), Inches(2.48), Inches(3.6), Inches(0.3),
     [("One Console AI — 자동 반영 루프", 12, True, BLUE)])
new = ["데이터 변경 감지 → Impact Analysis", "AI Ready Data 자동 갱신 · Vector Update",
       "Agent 자동 반영", "Feedback · Evaluation", "Prompt 개선 · Version 관리", "자동 배포 ↺"]
for i, t_ in enumerate(new):
    y = Inches(2.92) + Inches(0.56) * i
    dot(s, ML + Inches(4.05), y + Inches(0.12), Inches(0.16),
        BLUE if i in (0, 5) else WHITE, None if i in (0, 5) else HAIR)
    text(s, ML + Inches(4.3), y, Inches(3.55), Inches(0.45),
         [(t_, 10, i in (0, 5), INK if i in (0, 5) else MUT80)])
    if i < 5: conn(s, ML + Inches(4.05), y + Inches(0.2), ML + Inches(4.05), y + Inches(0.68), HAIR, 1.1)
# 우측: 고객 효과 6
ex = ML + Inches(8.3)
text(s, ex, Inches(2.35), Inches(3.6), Inches(0.3), [("고객 효과", 12, True, INK)])
effects = ["AI 구축 기간 단축", "운영 비용 절감", "AI 품질 지속 향상",
           "변경 영향 자동 분석", "AI 운영 표준화", "데이터 기반 의사결정"]
for i, t_ in enumerate(effects):
    y = Inches(2.78) + Inches(0.59) * i
    box(s, ex, y, Inches(3.6), Inches(0.46), PEARL, HAIR)
    text(s, ex + Inches(0.22), y + Inches(0.02), Inches(3.2), Inches(0.4),
         [([("✓ ", 10.5, True, BLUE), (t_, 10.5, True, INK)],)], anchor=MSO_ANCHOR.MIDDLE)
text(s, ML, Inches(6.55), CW, Inches(0.35),
     [("AI Ready Data 구축에서 끝나지 않고 — 데이터 변경부터 자동 배포까지, 운영 전 과정이 연결됩니다.", 12, True, INK)],
     align=PP_ALIGN.CENTER)

# ============================================================
# 13 · 사업화 (04) — Phase 1~3 + BM
# ============================================================
s = content("04", "실질적 사업화 추진",
            lead=lead_runs([("구축에서 운영까지, ", False), ("반복 매출 구조", True),
                            (" — MVP 기반 즉시 착수", False)]))
ph = [
    ("Phase 1", "Enterprise SI", "고객 맞춤 구축\nAI Ready Platform 딜리버리", True),
    ("Phase 2", "One Console AI SaaS", "AI Ready Data · AgentOps ·\nOntology · Prompt · Workflow 통합 관리", False),
    ("Phase 3", "AI MSP", "운영 · 평가 · 개선 · 모니터링\n지속 지원 (운영 반복 매출)", False),
]
for i, (p_, h_, b_, now) in enumerate(ph):
    x = ML + Inches(4.1) * i
    box(s, x, Inches(2.3), Inches(3.9), Inches(1.6), WHITE if now else PEARL,
        BLUE if now else HAIR, lw=1.4 if now else 1.0)
    text(s, x + Inches(0.28), Inches(2.48), Inches(3.35), Inches(1.25), [
        (p_, 9.5, True, BLUE, 3), (h_, 14, True, INK, 5), (b_, 9.5, False, MUT80)], spacing=1.18)
    if i < 2:
        text(s, x + Inches(3.9), Inches(2.9), Inches(0.2), Inches(0.4), [("→", 13, True, GRAY)],
             align=PP_ALIGN.CENTER)
text(s, ML, Inches(4.2), Inches(5), Inches(0.3), [("Business Model", 12.5, True, INK)])
bm = [("구축", "AI Ready Platform 구축 (SI)"),
      ("구독", "One Console AI 라이선스 (SaaS)"),
      ("운영", "MSP 운영 서비스 — 평가·개선 대행"),
      ("확장", "GraphRAG · Ontology · AI Governance")]
for i, (k, v) in enumerate(bm):
    y = Inches(4.6) + Inches(0.56) * i
    pill(s, ML, y, Inches(1.1), Inches(0.42), k, TILE if i < 3 else WHITE,
         WHITE if i < 3 else BLUE, 10, line=None if i < 3 else BLUE)
    text(s, ML + Inches(1.35), y + Inches(0.05), Inches(4.6), Inches(0.35),
         [(v, 10.5, False, MUT80)])
box(s, ML + Inches(6.5), Inches(4.2), Inches(5.6), Inches(2.5), TILE)
text(s, ML + Inches(6.85), Inches(4.45), Inches(4.95), Inches(2.05), [
    ("지식과 형상이 쌓일수록, 떠나기 어려워진다", 14, True, WHITE, 7),
    ("고객이 플랫폼에 축적한 AI Ready Data · Prompt · Ontology가 곧 전환 비용 — 1회성 SI가 3~5년 락인되는 구독·운영 매출로 바뀝니다.",
     10.5, False, BODY_D)], spacing=1.3)
text(s, ML, Inches(6.95), Inches(8), Inches(0.26),
     [("※ 가격·구성은 파일럿에서 지불의사 검증 후 확정", 8.5, False, GRAY)])

# ============================================================
# 14 · 로드맵 (04)
# ============================================================
s = content("04", "Roadmap",
            lead=lead_runs([("완벽한 계획보다, ", False), ("실행 가능한 첫걸음", True), (".", False)]))
ty = Inches(3.45)
conn(s, ML + Inches(0.5), ty, ML + Inches(11.6), ty, HAIR, 1.4)
phases = [
    ("현재", "MVP 검증 완료", ["콘솔 8화면 + API 19개 구현", "현장에서 운영 모델 확인"], "데모 가능", True),
    ("6개월", "그룹사 Pilot", ["유상 PoC 1~2사", "지불의사 · 가격 검증"], "Pilot 계약", False),
    ("12개월", "v1.0 제품 출시", ["Assessment · Data Studio 완성", "하이브리드 검색(Vector+Graph)"], "레퍼런스 확보", False),
    ("24개월", "SaaS 전환", ["멀티테넌트 · 구독 과금", "Ontology 자동 구축·영향 분석"], "반복 매출 본격화", False),
    ("36개월", "고객 확대", ["산업별 템플릿 · 파트너 채널", "AI MSP 운영 매출"], "고객 20~30사", False),
]
step_w = Inches(11.1) / 4
for i, (period, h_, bullets, kpi, active) in enumerate(phases):
    x = ML + Inches(0.5) + step_w * i
    dot(s, x, ty, Inches(0.26), BLUE if active else WHITE, None if active else HAIR)
    text(s, x - Inches(1.2), ty - Inches(0.95), Inches(2.4), Inches(0.72), [
        (period, 10.5, True, BLUE if active else MUT48, 3), (h_, 13, True, INK)],
        align=PP_ALIGN.CENTER)
    text(s, x - Inches(1.3), ty + Inches(0.28), Inches(2.6), Inches(0.95),
         [("· " + b_, 8.8, False, MUT80, 4) for b_ in bullets], align=PP_ALIGN.CENTER, spacing=1.15)
    pill(s, x - Inches(1.0), ty + Inches(1.42), Inches(2.0), Inches(0.38), kpi,
         BLUE if active else PEARL, WHITE if active else MUT80, 9,
         line=None if active else HAIR)
box(s, ML, Inches(5.75), CW, Inches(0.9), PEARL, HAIR)
text(s, ML + Inches(0.45), Inches(5.96), CW - Inches(0.9), Inches(0.55), [
    ([("첫걸음은 이미 뗐습니다 — ", 12.5, True, BLUE),
      ("필요한 것은 그룹사 Pilot 1개사와 소규모 TF. 6개월 내 유상 레퍼런스로 증명하겠습니다.", 12.5, True, INK)],)])

# ============================================================
# 15 · 신성장 동력 (05)
# ============================================================
s = content("05", "아이티센 신성장 동력",
            lead=lead_runs([("AI 프로젝트를 ", False), ("반복 가능한 플랫폼 사업", True), ("으로 전환.", False)]))
text(s, ML, Inches(2.25), Inches(6.2), Inches(0.35),
     [("프로젝트마다 소멸하던 것들이, 플랫폼의 디지털 자산으로 축적됩니다", 12, True, INK)])
assets = ["AI Ready Data", "Prompt", "Workflow", "Ontology", "Agent 설정", "운영 로그", "평가 결과"]
ax, ay = ML, Inches(2.75)
for i, a_ in enumerate(assets):
    x = ax + Inches(1.62) * (i % 4)
    y = ay + Inches(0.56) * (i // 4)
    pill(s, x, y, Inches(1.5), Inches(0.42), a_, PEARL, INK, 9, line=HAIR)
# 전환 다이어그램
box(s, ML, Inches(4.3), Inches(2.7), Inches(1.1), PEARL, HAIR)
text(s, ML, Inches(4.45), Inches(2.7), Inches(0.8), [
    ("기존", 9.5, True, MUT48, 3), ("SI 중심", 15, True, MUT48)],
    align=PP_ALIGN.CENTER)
text(s, ML + Inches(2.75), Inches(4.62), Inches(0.5), Inches(0.4), [("→", 16, True, GRAY)],
     align=PP_ALIGN.CENTER)
box(s, ML + Inches(3.3), Inches(4.3), Inches(2.9), Inches(1.1), WHITE, BLUE, lw=1.4)
text(s, ML + Inches(3.3), Inches(4.45), Inches(2.9), Inches(0.8), [
    ("미래", 9.5, True, BLUE, 3), ("Platform + SaaS + MSP", 13, True, BLUE)],
    align=PP_ALIGN.CENTER)
# 그룹 효과
gx = ML + Inches(6.7)
text(s, gx, Inches(2.25), Inches(5.2), Inches(0.35),
     [("그룹사 AI 프로젝트의 공통 플랫폼으로", 12, True, INK)])
geff = [("개발 생산성 향상", "표준 파이프라인·형상 재사용"),
        ("구축 기간 단축", "8주+ → 수일 수준"),
        ("품질 표준화", "전 프로젝트 동일 규칙·거버넌스"),
        ("운영 비용 절감", "중복 제거 + 자동 개선 루프")]
for i, (h_, b_) in enumerate(geff):
    y = Inches(2.72) + Inches(0.72) * i
    box(s, gx, y, Inches(5.4), Inches(0.6), PEARL, HAIR)
    text(s, gx + Inches(0.25), y + Inches(0.06), Inches(5.0), Inches(0.48), [
        ([("✓ ", 11, True, BLUE), (h_, 11.5, True, INK), ("   " + b_, 9.5, False, MUT48)],)],
        anchor=MSO_ANCHOR.MIDDLE)
box(s, ML, Inches(5.85), CW, Inches(0.92), TILE)
text(s, ML + Inches(0.45), Inches(6.02), CW - Inches(0.9), Inches(0.62), [
    ([("프로젝트가 끝나도 자산은 남는다 — ", 13, True, BLUE_D),
      ("수행할수록 강해지는 플랫폼이 아이티센의 새 성장 곡선이 됩니다.", 13, True, WHITE)],)],
    spacing=1.25)

# ============================================================
# 16 · 혁신 문화 & 구성원 성장 (06)
# ============================================================
s = content("06", "혁신 문화 & 구성원 성장",
            lead=lead_runs([("현장의 문제를 ", False), ("플랫폼 사업으로", True), (" 발전시키다.", False)]))
text(s, ML, Inches(2.3), CW, Inches(0.6), [
    ([("이 제안은 시장 조사에서 시작한 아이디어가 아니라, ", 12, False, MUT80),
      ("실제 고객 프로젝트에서 반복 경험한 문제를 해결하기 위해 직접 개발한 플랫폼", 12, True, INK),
      ("입니다.", 12, False, MUT80)],)], spacing=1.3)
j5 = ["문제 발굴", "서비스 기획", "MVP 개발", "제품화", "사업화"]
gap = Inches(0.07); fw = (CW - gap * 4) / 5
for i, t_ in enumerate(j5):
    done = i <= 2
    sp = box(s, ML + (fw + gap) * i, Inches(3.05), fw, Inches(0.72),
             BLUE if done else PEARL, None if done else HAIR,
             shape=MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t_; _set_font(r, 11, True, WHITE if done else MUT48)
text(s, ML, Inches(3.9), CW, Inches(0.3),
     [("구성원이 전 과정을 직접 경험 — 파란 단계는 이번 제안에서 이미 완료", 9.5, False, MUT48)])
grow = [("기술 자산", "전처리 파이프라인 · AgentOps 헬퍼 · UI 컴포넌트가 재사용 가능한 코드로"),
        ("사업 자산", "가격 모델 · 사업화 로드맵 · 고객 Pain Point 맵이 조직 지식으로"),
        ("조직 역량", "\"현장의 문제는 사업이 된다\"는 선례 — 도전이 장려되는 문화의 증거")]
for i, (h_, b_) in enumerate(grow):
    x = ML + Inches(4.1) * i
    kw_block(s, x, Inches(4.5), Inches(3.85), h_, b_)
box(s, ML, Inches(5.95), CW, Inches(0.85), PEARL, HAIR)
text(s, ML + Inches(0.45), Inches(6.14), CW - Inches(0.9), Inches(0.5), [
    ([("프로젝트 경험이 제품으로 축적되는 선순환 — ", 12.5, True, BLUE),
      ("기술·사업·조직이 함께 성장합니다.", 12.5, True, INK)],)])

# ============================================================
# 17 · Positioning (★) — 전체 생명주기 + One Message
# ============================================================
s = content("★", "One Console AI Positioning",
            lead=lead_runs([("AI 적용 대상 선정부터 지속 고도화까지 — ", False),
                            ("전 생명주기, 하나의 플랫폼", True), (".", False)]))
row1 = [("AI 적용 대상\n선정", False), ("AI Ready\nAssessment", False),
        ("AI Ready Data\n구축", True), ("Knowledge\nIntelligence", False), ("AI Agent\n생성", True)]
row2 = [("AgentOps", True), ("Continuous\nAI Operation", True),
        ("Feedback &\nEvaluation", True), ("지속적인\nAI 고도화", True)]
gap = Inches(0.14); fw1 = (CW - gap * 4) / 5
for i, (t_, done) in enumerate(row1):
    x = ML + (fw1 + gap) * i
    box(s, x, Inches(2.4), fw1, Inches(0.95), WHITE if done else PEARL,
        BLUE if done else HAIR, lw=1.4 if done else 1.0)
    text(s, x, Inches(2.46), fw1, Inches(0.84), [(t_, 10.5, True, INK)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.08)
    if i < 4:
        text(s, x + fw1 - Inches(0.02), Inches(2.68), Inches(0.2), Inches(0.35),
             [("→", 11, True, GRAY)], align=PP_ALIGN.CENTER)
# 우측 끝에서 아래로 꺾여 2행으로
conn(s, SW - ML - Inches(0.35), Inches(3.35), SW - ML - Inches(0.35), Inches(3.75), GRAY, 1.3)
text(s, SW - ML - Inches(0.55), Inches(3.42), Inches(0.4), Inches(0.3), [("↓", 11, True, GRAY)],
     align=PP_ALIGN.CENTER)
fw2 = (CW - gap * 3) / 4
for i, (t_, done) in enumerate(row2):
    x = SW - ML - (fw2 + gap) * (i + 1) + gap   # 오른쪽→왼쪽 흐름
    box(s, x, Inches(3.8), fw2, Inches(0.95), WHITE if done else PEARL,
        BLUE if done else HAIR, lw=1.4 if done else 1.0)
    text(s, x, Inches(3.86), fw2, Inches(0.84), [(t_, 11, True, INK)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.08)
    if i < 3:
        text(s, x - Inches(0.2), Inches(4.08), Inches(0.24), Inches(0.35),
             [("←", 11, True, GRAY)], align=PP_ALIGN.CENTER)
text(s, ML, Inches(4.95), CW, Inches(0.3),
     [("파란 테두리 = 현재 MVP로 동작 검증 완료 구간 · 회색 = v1.0 제품화 로드맵 구간", 9, False, MUT48)],
     align=PP_ALIGN.CENTER)
box(s, ML, Inches(5.45), CW, Inches(1.25), TILE)
text(s, ML + Inches(0.5), Inches(5.66), CW - Inches(1.0), Inches(0.9), [
    ([("One Message", 11, True, BLUE_D)], 6),
    ([("One Console AI는 AI Ready Data를 시작으로 AI Agent 구축·운영·개선까지 전 생명주기를 하나의 플랫폼에서 관리하는 ",
       13.5, True, WHITE),
      ("Enterprise AI AgentOps Platform", 13.5, True, BLUE_D), ("입니다.", 13.5, True, WHITE)],)],
    spacing=1.3)

# ============================================================
# 18 · 클로징
# ============================================================
s = base("black")
conn(s, ML, Inches(1.0), ML + Inches(0.44), Inches(1.0), BLUE, 2.4)
conn(s, ML + Inches(0.45), Inches(1.0), SW - ML, Inches(1.0), TILE2, 1.0)
text(s, ML, Inches(0.55), Inches(8), Inches(0.3),
     [("One Console AI — 마무리", 11, False, GRAY)])
text(s, ML, Inches(2.0), Inches(11.9), Inches(2.2), [
    ([("AI Ready Data", 42, True, BLUE_D, -1.0), ("에서 시작해,", 42, True, WHITE, -1.0)], 4),
    ([("AI Agent의 ", 42, True, WHITE, -1.0), ("전 생명주기", 42, True, BLUE_D, -1.0),
      ("로.", 42, True, WHITE, -1.0)],)], spacing=1.12)
box(s, ML + Inches(0.03), Inches(4.2), Inches(0.56), Inches(0.05), BLUE, r_px=99)
for i, (v, l) in enumerate([("5 STEP", "Assessment → Continuous Ops"),
                            ("MVP", "8개 화면 · 19개 API 동작 중"),
                            ("SI→SaaS→MSP", "반복 매출로의 전환")]):
    x = ML + Inches(4.0) * i
    text(s, x, Inches(4.6), Inches(3.7), Inches(1.1), [
        ([(v, 24, True, WHITE, -0.5)], 4), (l, 10.5, False, GRAY)])
box(s, 0, Inches(6.35), SW, Inches(0.62), TILE, shape=MSO_SHAPE.RECTANGLE)
text(s, ML, Inches(6.5), Inches(12.2), Inches(0.32),
     [("감사합니다   |   Enterprise AI AgentOps Platform   |   라이브 데모(콘솔 + Swagger UI) 준비 완료", 10.5, False, BODY_D)])

out = Path(__file__).resolve().parent / "One_Console_AI_제안서_v8.pptx"
prs.save(out)
print(f"saved: {out} / slides: {len(prs.slides._sldIdLst)}")
