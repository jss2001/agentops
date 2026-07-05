# -*- coding: utf-8 -*-
# One Console AI · 스플릿 표지 (좌 블루 그라디언트 패널 + 우 화이트 메시지)
# 실행: py proposal\make_cover.py → proposal\One_Console_AI_Cover.pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

BLUE   = RGBColor(0x00, 0x49, 0xF0)
BLUE_HI= RGBColor(0x1B, 0x4D, 0xF0)
INK    = RGBColor(0x0C, 0x0C, 0x0C)
MUT80  = RGBColor(0x2D, 0x2D, 0x2D)
GRAY   = RGBColor(0xA6, 0xA6, 0xA6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xBF, 0xD0, 0xFF)   # 좌측 패널 보조 텍스트
G1     = RGBColor(0x2E, 0x63, 0xFF)   # 그라디언트 밝은 쪽
G2     = RGBColor(0x0B, 0x2B, 0xD6)   # 그라디언트 어두운 쪽
FONT   = "Pretendard"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
s = prs.slides.add_slide(prs.slide_layouts[6])

def _set_font(run, size, bold, color, name=FONT, spc=None):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, name
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', name)
    if spc is not None:
        rPr.set('spc', str(int(spc * 100)))

def text(x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    tb = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        if len(para) == 5: p.space_after = Pt(para[4])
        if isinstance(para[0], list):
            for r_ in para[0]:
                run = p.add_run(); run.text = r_[0]
                _set_font(run, r_[1], r_[2], r_[3], spc=(r_[4] if len(r_) > 4 else None))
        else:
            run = p.add_run(); run.text = para[0]
            _set_font(run, para[1], para[2], para[3])
    return tb

def box(x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None, lw=1.0, r_px=14):
    sp = s.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = min(0.5, (r_px / 96) / min(w / 914400, h / 914400))
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp

# ---------- 배경 ----------
box(0, 0, SW, SH, WHITE)

# 좌측 블루 패널 (그라디언트 시도, 실패 시 솔리드)
panel = box(0, 0, Inches(5.33), SH, G2)
try:
    panel.fill.gradient()
    stops = panel.fill.gradient_stops
    stops[0].position = 0.0; stops[0].color.rgb = G1
    stops[1].position = 1.0; stops[1].color.rgb = G2
    panel.fill.gradient_angle = 125.0   # 좌상 밝음 → 우하 어두움
except Exception:
    panel.fill.solid(); panel.fill.fore_color.rgb = G2

# ---------- 좌측 패널 콘텐츠 ----------
box(Inches(0.75), Inches(1.78), Inches(0.58), Inches(0.032), WHITE, r_px=99,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(Inches(0.75), Inches(2.02), Inches(4.2), Inches(0.32),
     [([("ITCEN  ·  신사업 공모전", 12, True, WHITE, 2.2)],)])
text(Inches(0.72), Inches(2.55), Inches(4.4), Inches(2.6), [
    ([("One", 46, True, WHITE, -0.8)],),
    ([("Console", 46, True, WHITE, -0.8)],),
    ([("AI", 46, True, WHITE, -0.8)],)], spacing=1.04)
text(Inches(0.75), Inches(5.48), Inches(4.35), Inches(0.85), [
    ("Enterprise AI Operating System", 12.5, True, WHITE, 4),
    ("기업의 AI 지식과 운영 경험을 자산으로 축적하는 플랫폼", 10.5, False, LIGHT)],
    spacing=1.25)
box(Inches(0.78), Inches(6.95), Inches(0.09), Inches(0.09), WHITE, shape=MSO_SHAPE.OVAL)

# ---------- 우측 상단 로고 (아이티센엔텍) ----------
lx = Inches(11.0)
mark = box(lx, Inches(0.26), Inches(0.34), Inches(0.34), BLUE,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, r_px=8)
sail = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                          int(lx + Inches(0.075)), int(Inches(0.335)),
                          int(Inches(0.19)), int(Inches(0.16)))
sail.shadow.inherit = False
sail.fill.solid(); sail.fill.fore_color.rgb = WHITE; sail.line.fill.background()
text(lx + Inches(0.44), Inches(0.26), Inches(2.0), Inches(0.36),
     [([("아이티센", 15, True, INK), ("엔텍", 15, True, BLUE)],)], anchor=MSO_ANCHOR.MIDDLE)

# ---------- 우측 메시지 ----------
RX = Inches(6.19)
text(RX, Inches(1.78), Inches(5.5), Inches(0.32),
     [([("PROPOSAL 2026", 12, True, BLUE, 3.0)],)])
text(RX, Inches(2.3), Inches(6.6), Inches(2.7), [
    ([("AI Agent", 29, True, INK, -0.5), ("를 만드는 것이", 29, True, INK, -0.5)],),
    ([("문제가 아니라,", 29, True, INK, -0.5)],),
    ([("매번 다시 만드는 구조", 29, True, BLUE_HI, -0.5), ("가", 29, True, INK, -0.5)],),
    ([("문제입니다.", 29, True, INK, -0.5)],)], spacing=1.18)
box(RX, Inches(5.06), Inches(0.55), Inches(0.035), BLUE, r_px=99,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(RX, Inches(5.32), Inches(6.5), Inches(0.85), [
    ("AI 프로젝트가 끝나는 순간 지식과 운영 경험이 함께 사라지는", 11, False, MUT80, 3),
    ([("SI 중심의 일회성 사업을 ", 11, False, MUT80),
      ("“운영할수록 똑똑해지는 플랫폼”", 11, True, INK),
      ("으로 전환합니다.", 11, False, MUT80)],)], spacing=1.35)
text(RX, Inches(6.72), Inches(6.0), Inches(0.3),
     [([("ITCEN · Enterprise AI Platform 신사업 제안", 9.5, True, GRAY, 0.5)],)])

out = Path(__file__).resolve().parent / "One_Console_AI_Cover.pptx"
prs.save(out)
print(f"saved: {out} / slides: {len(prs.slides._sldIdLst)}")
